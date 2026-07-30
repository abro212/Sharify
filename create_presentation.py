import sys
import os
import qrcode
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

def build_powerpoint():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]

    # Color Palette
    EMERALD = RGBColor(6, 78, 59)      # #064E3B
    DARK_EMERALD = RGBColor(4, 62, 47) # #043E2F
    GOLD = RGBColor(197, 160, 89)      # #C5A059
    LIGHT_GOLD = RGBColor(245, 230, 200)
    CREAM = RGBColor(251, 249, 245)     # #FBF9F5
    DARK_TEXT = RGBColor(30, 41, 59)   # #1E293B
    MUTED_TEXT = RGBColor(100, 116, 139)
    WHITE = RGBColor(255, 255, 255)
    CARD_BG = RGBColor(241, 245, 249)

    def set_bg(slide, color):
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = color

    def add_header(slide, title_text, subtitle_text=None):
        tx_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(11.733), Inches(1.2))
        tf = tx_box.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0
        
        p = tf.paragraphs[0]
        p.text = title_text
        p.font.name = 'Georgia'
        p.font.size = Pt(28)
        p.font.bold = True
        p.font.color.rgb = EMERALD
        
        if subtitle_text:
            p2 = tf.add_paragraph()
            p2.text = subtitle_text
            p2.font.name = 'Arial'
            p2.font.size = Pt(14)
            p2.font.color.rgb = MUTED_TEXT
            p2.space_before = Pt(4)

    # -------------------------------------------------------------
    # SLIDE 1: Title Slide
    # -------------------------------------------------------------
    s1 = prs.slides.add_slide(blank_layout)
    set_bg(s1, CREAM)

    # Decorative Card Box
    card = s1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.5), Inches(1.2), Inches(10.333), Inches(5.1))
    card.fill.solid()
    card.fill.fore_color.rgb = WHITE
    card.line.color.rgb = GOLD
    card.line.width = Pt(2)

    tf1 = card.text_frame
    tf1.word_wrap = True
    tf1.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf1.margin_left = tf1.margin_right = Inches(0.5)

    p = tf1.paragraphs[0]
    p.text = "Pengabdian Masyarakat"
    p.font.name = 'Arial'
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = GOLD
    p.alignment = PP_ALIGN.CENTER

    p = tf1.add_paragraph()
    p.text = "Gerakan Edukasi Keuangan Syariah"
    p.font.name = 'Georgia'
    p.font.size = Pt(22)
    p.font.color.rgb = EMERALD
    p.alignment = PP_ALIGN.CENTER
    p.space_before = Pt(10)

    p = tf1.add_paragraph()
    p.text = "Smart Syariah Financial Planning"
    p.font.name = 'Georgia'
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = EMERALD
    p.alignment = PP_ALIGN.CENTER
    p.space_before = Pt(6)

    p = tf1.add_paragraph()
    p.text = "31 JULY 2026 | KAMPUS POLTEKKES KEMENKES JAKARTA 2"
    p.font.name = 'Arial'
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = GOLD
    p.alignment = PP_ALIGN.CENTER
    p.space_before = Pt(20)

    p = tf1.add_paragraph()
    p.text = "Program Studi Magister Manajemen Universitas Paramadina"
    p.font.name = 'Arial'
    p.font.size = Pt(14)
    p.font.color.rgb = DARK_TEXT
    p.alignment = PP_ALIGN.CENTER
    p.space_before = Pt(6)

    # -------------------------------------------------------------
    # SLIDE 2: Topik Pembahasan
    # -------------------------------------------------------------
    s2 = prs.slides.add_slide(blank_layout)
    set_bg(s2, CREAM)
    add_header(s2, "Topik Pembahasan", "Tiga Pilar Utama Pembekalan Literasi Keuangan Syariah")

    topics = [
      ("1. Prinsip Keuangan Syariah", "Memahami fondasi Fiqh Muamalah, Tauhid, Keadilan, & Keberkahan dalam harta."),
      ("2. Syariah Financial Planning", "Strategi formula 40-30-20-10, Zakat, Tabungan, Investasi Halal, & Dana Darurat."),
      ("3. Waspada Pinjol & Judol", "Mengidentifikasi bahaya Riba, Gharar, Maysir, serta pencegahan judi online.")
    ]

    for i, (top_title, top_desc) in enumerate(topics):
        left_pos = Inches(0.8 + i * 3.95)
        c = s2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left_pos, Inches(2.2), Inches(3.7), Inches(4.3))
        c.fill.solid()
        c.fill.fore_color.rgb = WHITE
        c.line.color.rgb = GOLD
        c.line.width = Pt(1.5)

        tf = c.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_right = Inches(0.3)
        tf.margin_top = Inches(0.4)

        p = tf.paragraphs[0]
        p.text = top_title
        p.font.name = 'Georgia'
        p.font.size = Pt(20)
        p.font.bold = True
        p.font.color.rgb = EMERALD

        p2 = tf.add_paragraph()
        p2.text = top_desc
        p2.font.name = 'Arial'
        p2.font.size = Pt(14)
        p2.font.color.rgb = DARK_TEXT
        p2.space_before = Pt(14)

    # -------------------------------------------------------------
    # SLIDE 3: Prinsip Keuangan Syariah
    # -------------------------------------------------------------
    s3 = prs.slides.add_slide(blank_layout)
    set_bg(s3, CREAM)
    add_header(s3, "Prinsip Keuangan Syariah", "Mengapa Islam Mengatur Keuangan?")

    # Big quote box
    qbox = s3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.8), Inches(11.733), Inches(1.3))
    qbox.fill.solid()
    qbox.fill.fore_color.rgb = EMERALD
    qbox.line.fill.background()
    tfq = qbox.text_frame
    tfq.word_wrap = True
    tfq.margin_left = Inches(0.4)
    p = tfq.paragraphs[0]
    p.text = "Uang bukan tujuan, melainkan alat untuk mencapai maslahat. Islam mengatur keuangan agar harta menjadi jalan kebaikan bukan sumber kerusakan."
    p.font.name = 'Georgia'
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = WHITE

    # 6 Principles Grid
    prinsip_list = [
        ("TAUHID", "Harta adalah milik Allah"),
        ("KEADILAN", "Keseimbangan hak & kewajiban"),
        ("TRANSPARANSI", "Akad jelas tanpa penipuan"),
        ("TIDAK MERUGIKAN", "Bebas dari bahaya dharar"),
        ("HALAL", "Bebas dari Riba, Gharar, Maysir"),
        ("KEBERKAHAN", "Berorientasi dunia & akhirat")
    ]

    for idx, (p_title, p_sub) in enumerate(prinsip_list):
        row = idx // 3
        col = idx % 3
        cx = Inches(0.8 + col * 3.95)
        cy = Inches(3.4 + row * 1.8)
        
        box = s3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, cx, cy, Inches(3.7), Inches(1.5))
        box.fill.solid()
        box.fill.fore_color.rgb = WHITE
        box.line.color.rgb = GOLD
        box.line.width = Pt(1)

        tf = box.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0.25)
        tf.margin_top = Inches(0.2)

        p = tf.paragraphs[0]
        p.text = p_title
        p.font.name = 'Georgia'
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = EMERALD

        p2 = tf.add_paragraph()
        p2.text = p_sub
        p2.font.name = 'Arial'
        p2.font.size = Pt(12)
        p2.font.color.rgb = MUTED_TEXT
        p2.space_before = Pt(4)

    # -------------------------------------------------------------
    # SLIDE 4: Apa itu Riba?
    # -------------------------------------------------------------
    s4 = prs.slides.add_slide(blank_layout)
    set_bg(s4, CREAM)
    add_header(s4, "Apa itu Riba?", "Tambahan yang dipersyaratkan dalam pertukaran atau pinjaman tanpa imbalan usaha/risiko sepadan.")

    r1 = s4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(2.0), Inches(5.7), Inches(4.7))
    r1.fill.solid()
    r1.fill.fore_color.rgb = WHITE
    r1.line.color.rgb = GOLD
    r1.line.width = Pt(1.5)

    tf = r1.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Inches(0.4)
    tf.margin_top = Inches(0.4)
    p = tf.paragraphs[0]
    p.text = "Riba Fadhl"
    p.font.name = 'Georgia'
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = EMERALD

    p2 = tf.add_paragraph()
    p2.text = "Tambahan pada tukar-menukar barang ribawi sejenis dengan takaran/timbangan yang berbeda.\n\nContoh: Menukar emas 10 gram dengan emas 11 gram; menukar uang Rp100 ribu dengan pecahan baru senilai Rp95 ribu."
    p2.font.name = 'Arial'
    p2.font.size = Pt(14)
    p2.font.color.rgb = DARK_TEXT
    p2.space_before = Pt(10)

    r2 = s4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.8), Inches(2.0), Inches(5.7), Inches(4.7))
    r2.fill.solid()
    r2.fill.fore_color.rgb = WHITE
    r2.line.color.rgb = GOLD
    r2.line.width = Pt(1.5)

    tf2 = r2.text_frame
    tf2.word_wrap = True
    tf2.margin_left = tf2.margin_right = Inches(0.4)
    tf2.margin_top = Inches(0.4)
    p = tf2.paragraphs[0]
    p.text = "Riba Nasi'ah"
    p.font.name = 'Georgia'
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = EMERALD

    p2 = tf2.add_paragraph()
    p2.text = "Tambahan yang dipersyaratkan karena penundaan waktu pembayaran utang-piutang.\n\nContoh: Pinjam Rp1 juta wajib kembali Rp1.2 juta; bunga kartu kredit & denda pinjol ilegal yang terus berbunga majemuk."
    p2.font.name = 'Arial'
    p2.font.size = Pt(14)
    p2.font.color.rgb = DARK_TEXT
    p2.space_before = Pt(10)

    # -------------------------------------------------------------
    # SLIDE 5: Survei Membuktikan
    # -------------------------------------------------------------
    s5 = prs.slides.add_slide(blank_layout)
    set_bg(s5, CREAM)
    add_header(s5, "Survei Membuktikan", "Sumber Data: OJK, BPS SNLIK 2024, Satgas PASTI 2026, PPATK 2025-2026")

    stats = [
        ("39,11%", "Indeks literasi keuangan syariah nasional (masih perlu ditingkatkan)."),
        ("951", "Entitas pinjol ilegal yang telah diblokir OJK."),
        ("12,3 jt", "Pemain judi online yang terdeteksi di Indonesia."),
        ("80%", "Perputaran dana judi online berasal dari pelajar & mahasiswa.")
    ]

    for i, (num, desc) in enumerate(stats):
        left_pos = Inches(0.8 + i * 2.95)
        box = s5.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left_pos, Inches(2.2), Inches(2.7), Inches(4.3))
        box.fill.solid()
        box.fill.fore_color.rgb = EMERALD
        box.line.color.rgb = GOLD
        box.line.width = Pt(2)

        tf = box.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_right = Inches(0.2)
        tf.margin_top = Inches(0.5)

        p = tf.paragraphs[0]
        p.text = num
        p.font.name = 'Georgia'
        p.font.size = Pt(36)
        p.font.bold = True
        p.font.color.rgb = GOLD
        p.alignment = PP_ALIGN.CENTER

        p2 = tf.add_paragraph()
        p2.text = desc
        p2.font.name = 'Arial'
        p2.font.size = Pt(13)
        p2.font.color.rgb = WHITE
        p2.alignment = PP_ALIGN.CENTER
        p2.space_before = Pt(16)

    # -------------------------------------------------------------
    # SLIDE 6: Kenapa Mahasiswa Jadi Sasaran Empuk?
    # -------------------------------------------------------------
    s6 = prs.slides.add_slide(blank_layout)
    set_bg(s6, CREAM)
    add_header(s6, "Kenapa Mahasiswa Jadi Sasaran Empuk?", "Faktor Kerentanan Generasi Muda Terhadap Pinjol Ilegal & Judol")

    reasons = [
        ("Tekanan Sosial", "Uang saku terbatas, namun tekanan gaya hidup & nongkrong tinggi."),
        ("Digital Native", "Sangat familiar dengan aplikasi mobile & transaksi digital instan."),
        ("FOMO (Fear Of Missing Out)", "Terpapar iklan pinjol/judol yang agresif & takut ketinggalan tren."),
        ("Literasi Syariah Rendah", "Minim pemahaman akan bahaya serta larangan Riba, Gharar, & Maysir."),
        ("Butuh Dana Instan", "Tergiur proses pencairan dana cepat tanpa verifikasi yang rumit.")
    ]

    for i, (title, desc) in enumerate(reasons):
        top_pos = Inches(1.8 + i * 1.05)
        box = s6.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), top_pos, Inches(11.733), Inches(0.9))
        box.fill.solid()
        box.fill.fore_color.rgb = WHITE
        box.line.color.rgb = GOLD

        tf = box.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0.4)
        tf.margin_top = Inches(0.15)

        p = tf.paragraphs[0]
        p.text = f"{title}: "
        p.font.name = 'Georgia'
        p.font.size = Pt(15)
        p.font.bold = True
        p.font.color.rgb = EMERALD

        run = p.add_run()
        run.text = desc
        run.font.name = 'Arial'
        run.font.size = Pt(14)
        run.font.bold = False
        run.font.color.rgb = DARK_TEXT

    # -------------------------------------------------------------
    # SLIDE 7: 5 Langkah Memulai
    # -------------------------------------------------------------
    s7 = prs.slides.add_slide(blank_layout)
    set_bg(s7, CREAM)
    add_header(s7, "5 Langkah Memulai Smart Syariah Financial Planning", "Tahapan Perencanaan Keuangan Berkah untuk Mahasiswa & Keluarga")

    steps = [
        ("1. Mindset Halal", "Fokus mencari keberkahan harta"),
        ("2. Rencana Anggaran", "Formula 40-30-20-10"),
        ("3. Zakat / Sedekah", "Utamakan berbagi di awal"),
        ("4. Investasi Halal", "Reksadana & Emas Syariah"),
        ("5. Dana Darurat", "Minimal 3x pengeluaran")
    ]

    for i, (st_t, st_d) in enumerate(steps):
        left_pos = Inches(0.8 + i * 2.38)
        c = s7.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left_pos, Inches(2.3), Inches(2.2), Inches(4.2))
        c.fill.solid()
        c.fill.fore_color.rgb = WHITE
        c.line.color.rgb = GOLD
        c.line.width = Pt(1.5)

        tf = c.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_right = Inches(0.15)
        tf.margin_top = Inches(0.4)

        p = tf.paragraphs[0]
        p.text = st_t
        p.font.name = 'Georgia'
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = EMERALD
        p.alignment = PP_ALIGN.CENTER

        p2 = tf.add_paragraph()
        p2.text = st_d
        p2.font.name = 'Arial'
        p2.font.size = Pt(13)
        p2.font.color.rgb = DARK_TEXT
        p2.alignment = PP_ALIGN.CENTER
        p2.space_before = Pt(14)

    # -------------------------------------------------------------
    # SLIDE 8: Halal-Way is My Way
    # -------------------------------------------------------------
    s8 = prs.slides.add_slide(blank_layout)
    set_bg(s8, CREAM)
    add_header(s8, "Halal-Way is My Way", "Gaya Hidup Halal Menuju Ketenangan Finansial")

    c8 = s8.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.5), Inches(2.0), Inches(10.333), Inches(4.5))
    c8.fill.solid()
    c8.fill.fore_color.rgb = WHITE
    c8.line.color.rgb = GOLD
    c8.line.width = Pt(2)

    tf8 = c8.text_frame
    tf8.word_wrap = True
    tf8.margin_left = tf8.margin_right = Inches(0.6)
    tf8.margin_top = Inches(0.6)

    p = tf8.paragraphs[0]
    p.text = "Fokus Pada Keberkahan & Ketenangan Hidup"
    p.font.name = 'Georgia'
    p.font.size = Pt(26)
    p.font.bold = True
    p.font.color.rgb = EMERALD

    p2 = tf8.add_paragraph()
    p2.text = "Hindari Riba, Gharar, dan Maysir dalam segala bentuk transaksi keuangan harian.\n\nUbah mindset dari sekadar 'mencari untung sebanyak-banyaknya' menjadi 'fokus mencari keberkahan dalam hidup' yang akan membawa kita pada ketenangan, kedamaian, dan ridha Allah SWT."
    p2.font.name = 'Arial'
    p2.font.size = Pt(18)
    p2.font.color.rgb = DARK_TEXT
    p2.space_before = Pt(16)

    # -------------------------------------------------------------
    # SLIDE 9: My Smart Syariah Financial Planning (Formula)
    # -------------------------------------------------------------
    s9 = prs.slides.add_slide(blank_layout)
    set_bg(s9, CREAM)
    add_header(s9, "My Smart Syariah Financial Planning", "Buat Rencana Keuangan Syariahmu dengan Formula 40-30-20-10")

    f_box = s9.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(2.0), Inches(11.733), Inches(4.7))
    f_box.fill.solid()
    f_box.fill.fore_color.rgb = WHITE
    f_box.line.color.rgb = GOLD
    f_box.line.width = Pt(1.5)

    tf9 = f_box.text_frame
    tf9.word_wrap = True
    tf9.margin_left = tf9.margin_right = Inches(0.6)
    tf9.margin_top = Inches(0.4)

    p = tf9.paragraphs[0]
    p.text = "Formula Alokasi Keuangan Ideal:"
    p.font.name = 'Georgia'
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = EMERALD

    p2 = tf9.add_paragraph()
    p2.text = "• 40% Kebutuhan Pokok (Makan, tempat tinggal, kewajiban harian)\n• 30% Pendidikan & Transportasi (Uang kuliah, buku, bensin/ongkos)\n• 20% Tabungan & Investasi Halal (Reksadana syariah, emas digital)\n• 10% Sedekah & Zakat (Berbagi untuk sesama di awal gajian)"
    p2.font.name = 'Arial'
    p2.font.size = Pt(16)
    p2.font.color.rgb = DARK_TEXT
    p2.space_before = Pt(12)

    p3 = tf9.add_paragraph()
    p3.text = "Yuk, coba buat rencana keuanganmu sendiri melalui tautan:\nhttps://s.id/mysmartsyariahfinancialplanning"
    p3.font.name = 'Arial'
    p3.font.size = Pt(15)
    p3.font.bold = True
    p3.font.color.rgb = GOLD
    p3.space_before = Pt(20)

    # -------------------------------------------------------------
    # SLIDE 10: Tabungan & Investasi
    # -------------------------------------------------------------
    s10 = prs.slides.add_slide(blank_layout)
    set_bg(s10, CREAM)
    add_header(s10, "Tabungan & Investasi", "Pilih Instrumen Investasi Syariah yang Aman & Terdaftar")

    box10 = s10.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(2.0), Inches(11.733), Inches(4.7))
    box10.fill.solid()
    box10.fill.fore_color.rgb = WHITE
    box10.line.color.rgb = GOLD
    box10.line.width = Pt(1.5)

    tf10 = box10.text_frame
    tf10.word_wrap = True
    tf10.margin_left = tf10.margin_right = Inches(0.5)
    tf10.margin_top = Inches(0.4)

    p = tf10.paragraphs[0]
    p.text = "Instrumen Investasi Halal Terdaftar OJK & DSN-MUI:"
    p.font.name = 'Georgia'
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = EMERALD

    p2 = tf10.add_paragraph()
    p2.text = "Mulai berinvestasi pada Reksadana Syariah, Emas Digital, Saham Syariah (JII & ISSI), serta Sukuk Negara.\n\nPastikan platform sekuritas/aplikasi terdaftar resmi dan berizin OJK, serta mengantongi sertifikasi kepatuhan Fatwa DSN-MUI."
    p2.font.name = 'Arial'
    p2.font.size = Pt(16)
    p2.font.color.rgb = DARK_TEXT
    p2.space_before = Pt(12)

    p3 = tf10.add_paragraph()
    p3.text = "💡 Trik Investasi Mahasiswa:\nGunakan metode investasi rutin per bulan (Dollar-Cost Averaging / DCA) mulai dari nominal kecil Rp10.000 secara konsisten."
    p3.font.name = 'Arial'
    p3.font.size = Pt(15)
    p3.font.bold = True
    p3.font.color.rgb = GOLD
    p3.space_before = Pt(16)

    # -------------------------------------------------------------
    # SLIDE 11: Dana Darurat
    # -------------------------------------------------------------
    s11 = prs.slides.add_slide(blank_layout)
    set_bg(s11, CREAM)
    add_header(s11, "Dana Darurat", "Fondasi Keamanan Finansial Sebelum Berinvestasi")

    box11 = s11.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.5), Inches(2.0), Inches(10.333), Inches(4.5))
    box11.fill.solid()
    box11.fill.fore_color.rgb = WHITE
    box11.line.color.rgb = GOLD
    box11.line.width = Pt(2)

    tf11 = box11.text_frame
    tf11.word_wrap = True
    tf11.margin_left = tf11.margin_right = Inches(0.6)
    tf11.margin_top = Inches(0.5)

    p = tf11.paragraphs[0]
    p.text = "Sisihkan Dana Darurat Minimal 3x Pengeluaran Bulanan"
    p.font.name = 'Georgia'
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = EMERALD

    p2 = tf11.add_paragraph()
    p2.text = "Dana darurat berfungsi sebagai benteng perlindungan utama agar saat terjadi musibah mendadak (sakit, ganti laptop, atau krisis), kita tidak terjebak menggunakan pinjol atau berutang riba.\n\nSimpan dana darurat di instrumen cair yang aman seperti Tabungan Syariah Wadi'ah atau Reksadana Pasar Uang Syariah."
    p2.font.name = 'Arial'
    p2.font.size = Pt(16)
    p2.font.color.rgb = DARK_TEXT
    p2.space_before = Pt(16)

    # -------------------------------------------------------------
    # SLIDE 12 (FEATURED NEW SLIDE): Sharify AI Co-Pilot & QR Code
    # -------------------------------------------------------------
    s12 = prs.slides.add_slide(blank_layout)
    set_bg(s12, CREAM)
    add_header(s12, "Solusi Praktis: Cobalah Sharify AI Co-Pilot!", "Asisten Keuangan Syariah Pintar Berbasis AI & Fatwa DSN-MUI dalam Satu Genggaman")

    # Left Feature Card
    fcard = s12.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.9), Inches(7.5), Inches(5.0))
    fcard.fill.solid()
    fcard.fill.fore_color.rgb = WHITE
    fcard.line.color.rgb = GOLD
    fcard.line.width = Pt(2)

    tff = fcard.text_frame
    tff.word_wrap = True
    tff.margin_left = tff.margin_right = Inches(0.4)
    tff.margin_top = Inches(0.3)

    p = tff.paragraphs[0]
    p.text = "Fitur Pintar Sharify App:"
    p.font.name = 'Georgia'
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = EMERALD

    features = [
        ("🤖 AI Chatbot Syariah", "Konsultasi hukum Fiqh Muamalah & Fatwa DSN-MUI 24/7 (Anti-Ngarang)."),
        ("🧮 Smart Zakat & Faraidh Calculator", "Hitung zakat profesi, saham, & waris otomatis secara presisi."),
        ("🛡️ Riba Detox & Halal Asset Screener", "Simulasi bebas utang riba & cek kepatuhan syariah aset Anda.")
    ]

    for f_title, f_desc in features:
        p1 = tff.add_paragraph()
        p1.text = f_title
        p1.font.name = 'Georgia'
        p1.font.size = Pt(15)
        p1.font.bold = True
        p1.font.color.rgb = EMERALD
        p1.space_before = Pt(10)

        p2 = tff.add_paragraph()
        p2.text = f_desc
        p2.font.name = 'Arial'
        p2.font.size = Pt(13)
        p2.font.color.rgb = DARK_TEXT

    # Right QR Code Card
    qrcard = s12.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(8.6), Inches(1.9), Inches(3.9), Inches(5.0))
    qrcard.fill.solid()
    qrcard.fill.fore_color.rgb = EMERALD
    qrcard.line.color.rgb = GOLD
    qrcard.line.width = Pt(2)

    tfqr = qrcard.text_frame
    tfqr.word_wrap = True
    tfqr.margin_left = tfqr.margin_right = Inches(0.2)
    tfqr.margin_top = Inches(0.3)

    p = tfqr.paragraphs[0]
    p.text = "SCAN UNTUK COBA"
    p.font.name = 'Georgia'
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = GOLD
    p.alignment = PP_ALIGN.CENTER

    # Add QR Code Image
    qr_img_path = 'public/sharify_qr_code.png'
    if os.path.exists(qr_img_path):
        s12.shapes.add_picture(qr_img_path, Inches(9.2), Inches(2.6), Inches(2.7), Inches(2.7))

    # Add URL Label at bottom of QR card
    tx_url = s12.shapes.add_textbox(Inches(8.7), Inches(5.5), Inches(3.7), Inches(1.2))
    tfu = tx_url.text_frame
    tfu.word_wrap = True
    p = tfu.paragraphs[0]
    p.text = "Akses Langsung via Browser:"
    p.font.name = 'Arial'
    p.font.size = Pt(11)
    p.font.color.rgb = LIGHT_GOLD
    p.alignment = PP_ALIGN.CENTER

    p2 = tfu.add_paragraph()
    p2.text = "s.id/mysmartsyariahfinancialplanning"
    p2.font.name = 'Arial'
    p2.font.size = Pt(12)
    p2.font.bold = True
    p2.font.color.rgb = WHITE
    p2.alignment = PP_ALIGN.CENTER
    p2.space_before = Pt(2)

    # -------------------------------------------------------------
    # SLIDE 13: Tonight's Treasures
    # -------------------------------------------------------------
    s13 = prs.slides.add_slide(blank_layout)
    set_bg(s13, CREAM)
    add_header(s13, "Tonight's Treasures", "Apresiasi Peserta Pengabdian Masyarakat")

    c13_1 = s13.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(2.0), Inches(5.7), Inches(4.7))
    c13_1.fill.solid()
    c13_1.fill.fore_color.rgb = WHITE
    c13_1.line.color.rgb = GOLD
    c13_1.line.width = Pt(1.5)

    tf = c13_1.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Inches(0.4)
    tf.margin_top = Inches(0.4)
    p = tf.paragraphs[0]
    p.text = "Grand Prize"
    p.font.name = 'Georgia'
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = EMERALD

    p2 = tf.add_paragraph()
    p2.text = "The winning team takes home the Studio Shodwe 'Royal Iftar' Hamper.\n\nIncludes rare honey, Zamzam water, premium Ajwa dates, and custom prayer mats."
    p2.font.name = 'Arial'
    p2.font.size = Pt(15)
    p2.font.color.rgb = DARK_TEXT
    p2.space_before = Pt(12)

    c13_2 = s13.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.8), Inches(2.0), Inches(5.7), Inches(4.7))
    c13_2.fill.solid()
    c13_2.fill.fore_color.rgb = WHITE
    c13_2.line.color.rgb = GOLD
    c13_2.line.width = Pt(1.5)

    tf2 = c13_2.text_frame
    tf2.word_wrap = True
    tf2.margin_left = tf2.margin_right = Inches(0.4)
    tf2.margin_top = Inches(0.4)
    p = tf2.paragraphs[0]
    p.text = "Runner Up"
    p.font.name = 'Georgia'
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = EMERALD

    p2 = tf2.add_paragraph()
    p2.text = "Gift vouchers for the finest Halal dining experience in the city.\n\nVoucher kuliner halal eksklusif & merchandise spesial."
    p2.font.name = 'Arial'
    p2.font.size = Pt(15)
    p2.font.color.rgb = DARK_TEXT
    p2.space_before = Pt(12)

    # -------------------------------------------------------------
    # SLIDE 14: Terima Kasih
    # -------------------------------------------------------------
    s14 = prs.slides.add_slide(blank_layout)
    set_bg(s14, CREAM)

    card14 = s14.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.5), Inches(1.0), Inches(10.333), Inches(5.5))
    card14.fill.solid()
    card14.fill.fore_color.rgb = EMERALD
    card14.line.color.rgb = GOLD
    card14.line.width = Pt(2)

    tf14 = card14.text_frame
    tf14.word_wrap = True
    tf14.margin_left = tf14.margin_right = Inches(0.5)
    tf14.margin_top = Inches(0.5)

    p = tf14.paragraphs[0]
    p.text = "Terima Kasih"
    p.font.name = 'Georgia'
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = GOLD
    p.alignment = PP_ALIGN.CENTER

    p2 = tf14.add_paragraph()
    p2.text = "KELOMPOK MAHASISWA MAGISTER MANAJEMEN\nUNIVERSITAS PARAMADINA"
    p2.font.name = 'Arial'
    p2.font.size = Pt(14)
    p2.font.bold = True
    p2.font.color.rgb = WHITE
    p2.alignment = PP_ALIGN.CENTER
    p2.space_before = Pt(10)

    p3 = tf14.add_paragraph()
    p3.text = "Asep Ajaeni  •  Ayu Nadia Faishol  •  Fachrul Rozie Yudha Gunawan\nGemma Santoni Harahap  •  Harimurti Mustikayudha Satriyatama\nIqbal Rizky Nugroho  •  Nova Indah T  •  Rahmat Abrori\nShinta Asrini  •  Tomi Margiantara Sugito"
    p3.font.name = 'Arial'
    p3.font.size = Pt(14)
    p3.font.color.rgb = LIGHT_GOLD
    p3.alignment = PP_ALIGN.CENTER
    p3.space_before = Pt(20)

    # Save presentation
    output_path = 'public/Smart_Syariah_Financial_Planning_Sharify.pptx'
    prs.save(output_path)
    print(f"Presentation saved successfully to: {output_path}")

if __name__ == '__main__':
    build_powerpoint()
